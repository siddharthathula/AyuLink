import sys
import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = '/home/siddhartha/Desktop/Template for Techsphere presentations.pptx.pptx'
OUTPUT_PATH_1 = TEMPLATE_PATH
OUTPUT_PATH_2 = '/home/siddhartha/Desktop/AyuLink_TechSphere_2026_Presentation.pptx'

# Colors
TEAL = RGBColor(13, 148, 136)       # #0d9488
DARK_SLATE = RGBColor(15, 23, 42)   # #0f172a
MUTED_SLATE = RGBColor(71, 85, 105) # #475569
EMERALD = RGBColor(16, 185, 129)    # #10b981
ROSE = RGBColor(225, 29, 72)        # #e11d48

def get_or_create_body_textbox(slide, left=Inches(0.8), top=Inches(1.8), width=Inches(6.8), height=Inches(5.0)):
    # Check if TextBox 9 or similar body box exists
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name in ['TextBox 9', 'TextBox 8'] and shape.name != slide.shapes[0].name:
            if 'May use' in shape.text or 'One slide' in shape.text or 'Follow similar' in shape.text or not shape.text.strip():
                return shape.text_frame
    # Otherwise create one
    txBox = slide.shapes.add_textbox(left, top, width, height)
    return txBox.text_frame

def populate_presentation():
    prs = pptx.Presentation(TEMPLATE_PATH)
    print(f"Loaded template with {len(prs.slides)} slides.")

    # ── SLIDE 1: Title Slide ──────────────────────────────────────────────────
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.has_text_frame and 'TextBox 8' in shape.name:
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = "AyuLink — When Every Second Counts"
            p0.font.bold = True
            p0.font.size = Pt(26)
            p0.font.color.rgb = TEAL

            p1 = tf.add_paragraph()
            p1.text = "Universal IoT & AI Healthcare Platform for Rural, Urban & Telemedicine Patients"
            p1.font.size = Pt(15)
            p1.font.color.rgb = MUTED_SLATE

            p2 = tf.add_paragraph()
            p2.text = "\nTeam FightClub: Siddhartha & Anirudh"
            p2.font.bold = True
            p2.font.size = Pt(17)
            p2.font.color.rgb = DARK_SLATE

            p3 = tf.add_paragraph()
            p3.text = "Woxsen University | TechSphere 2026"
            p3.font.size = Pt(13)
            p3.font.color.rgb = MUTED_SLATE

    # ── SLIDE 2: Table of Contents ────────────────────────────────────────────
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.has_text_frame and 'TextBox 8' in shape.name:
            tf = shape.text_frame
            tf.clear()
            toc_items = [
                "1. Introduction & Healthcare Emergency Gap",
                "2. Literature Review & Existing Solutions Gap",
                "3. Project Objectives",
                "4. System Methodology & Architecture (LoRa + FastAPI + AI)",
                "5. Project Implementation (Hardware + Backend + Next.js 16 Dashboard)",
                "6. Results & Discussion (Sub-3s Latency, 5km Mesh Coverage)",
                "7. Impact of Project on Society",
                "8. Future Scope & District Scalability Roadmap",
                "9. Conclusions",
                "10. References"
            ]
            for idx, item in enumerate(toc_items):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_SLATE
                p.font.bold = True

    # ── SLIDE 3: Introduction ────────────────────────────────────────────────
    s3 = prs.slides[2]
    tf3 = get_or_create_body_textbox(s3)
    tf3.clear()
    bullets3 = [
        ("The Reality in Rural India", "47 minutes average emergency response time in rural areas vs 8 minutes in urban centers."),
        ("The Unmonitored Population", "600M+ Indians above age 40 have zero continuous health monitoring."),
        ("Barriers in Existing Tech", "Existing wearables require smartphones, 4G/WiFi, and high monthly costs (₹500+/mo)."),
        ("The AyuLink Solution", "Custom ESP32-S3 wristband with MAX30102 HR/SpO2, MPU6050 fall/tremor detection, and 5km LoRa radio operating at ₹18/month."),
        ("Tri-Tier Healthcare Ecosystem", "Serves Rural (offline LoRa mesh), Urban (WiFi Smart Hub), and Telemedicine (Zero-hardware digital ABHA sync) seamlessly.")
    ]
    for idx, (title, desc) in enumerate(bullets3):
        p = tf3.paragraphs[0] if idx == 0 else tf3.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_SLATE

    img3_path = "/home/siddhartha/Desktop/copied/backfol/AyuLink/ElderCare/photos/wristband_oled_mpu6050.jpg"
    if os.path.exists(img3_path):
        s3.shapes.add_picture(img3_path, Inches(8.0), Inches(2.0), width=Inches(4.8))

    # ── SLIDE 4: Literature Review ───────────────────────────────────────────
    s4 = prs.slides[3]
    tf4 = get_or_create_body_textbox(s4, width=Inches(11.5))
    tf4.clear()
    bullets4 = [
        ("Existing RPM Wearables (Apple/Fitbit)", "High cost (₹15,000+), depend strictly on Bluetooth/WiFi/4G, fail completely in rural zero-connectivity zones."),
        ("Traditional Telemedicine Platforms", "Require active smartphone user interaction, making them useless during sudden cardiac events or unconscious falls."),
        ("Language & Literacy Barriers", "Most medical apps are English-only, excluding uneducated rural families and field workers."),
        ("AyuLink Key Innovation", "Hardware-agnostic 433MHz LoRa mesh, sub-3-second emergency alert latency, Groq LLaMA 3.1 clinical triage, and native Telugu/Hindi voice alerts.")
    ]
    for idx, (title, desc) in enumerate(bullets4):
        p = tf4.paragraphs[0] if idx == 0 else tf4.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE

    # ── SLIDE 5: Objectives ──────────────────────────────────────────────────
    s5 = prs.slides[4]
    tf5 = get_or_create_body_textbox(s5, width=Inches(11.5))
    tf5.clear()
    bullets5 = [
        ("Low-Cost Hardware", "Develop an affordable RPM wearable costing <₹2,500 ($30) with ₹18/month operating cost."),
        ("Offline LoRa Mesh", "Build a 5km 433MHz LoRa radio bridge to transmit vitals without WiFi or smartphones."),
        ("Real-Time AI Clinical Triage", "Implement Groq LLaMA 3.1 & local Ollama for sub-second medical triage and 0–100 risk scoring."),
        ("Modern Command Dashboard", "Build a Next.js 16 + React 19 dark mode glassmorphism UI with real-time WebSockets."),
        ("Multilingual Emergency Response", "Push Telugu/Hindi voice alerts via Telegram with instant 108 ambulance dispatch.")
    ]
    for idx, (title, desc) in enumerate(bullets5):
        p = tf5.paragraphs[0] if idx == 0 else tf5.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE

    # ── SLIDE 6: Methodology ─────────────────────────────────────────────────
    s6 = prs.slides[5]
    tf6 = get_or_create_body_textbox(s6)
    tf6.clear()
    bullets6 = [
        ("Physical Hardware Layer", "ESP32-S3 Wearable (MAX30102 + MPU6050 + LoRa SX1278 + SOS) → ESP32 Gateway (LoRa to WS bridge) → Smart Hub (4-slot servo pill dispenser + MQ-135 + DHT11 + Flame) → ESP32-CAM MJPEG stream."),
        ("Backend Processing Pipeline", "FastAPI Python server + SQLite WAL database + WebSocket duplex channels + ThresholdEngine clinical filter."),
        ("AI Triage & Natural Language", "Groq LLaMA 3.1 8B instant triage + mental health distress keyword intercepter + drug/disease search engine."),
        ("Multilingual Voice Notification", "Telemetry anomaly → AI Triage → gTTS Telugu/Hindi voice note → Telegram API push with 108 dial button.")
    ]
    for idx, (title, desc) in enumerate(bullets6):
        p = tf6.paragraphs[0] if idx == 0 else tf6.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_SLATE

    img6_path = "/home/siddhartha/Desktop/copied/backfol/AyuLink/ElderCare/photos/full_hardware_setup.jpg"
    if os.path.exists(img6_path):
        s6.shapes.add_picture(img6_path, Inches(7.8), Inches(2.0), width=Inches(5.0))

    # ── SLIDE 7: Project Implementation ──────────────────────────────────────
    s7 = prs.slides[6]
    tf7 = get_or_create_body_textbox(s7)
    tf7.clear()
    bullets7 = [
        ("Hardware Prototyping", "Hand-soldered ESP32-S3 node with MAX30102, MPU6050, SSD1306 OLED, and LoRa SX1278 antenna."),
        ("Backend & Alert Engine", "Async FastAPI server handling 100+ concurrent WebSockets, SQLite persistent storage, camera proxy loop, and ThresholdEngine anomaly detection."),
        ("AI Clinical Triage", "Groq LLaMA 3.1 8B instant triage engine, mental health distress keyword intercepter, and drug/disease reference search engine."),
        ("Web Dashboard & Family Portal", "Next.js 16 App Router interface featuring live vitals charts, emergency GPS maps, ASHA visit logging, and responsive mobile layouts.")
    ]
    for idx, (title, desc) in enumerate(bullets7):
        p = tf7.paragraphs[0] if idx == 0 else tf7.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_SLATE

    img7_path = "/home/siddhartha/Desktop/copied/backfol/AyuLink/ElderCare/screenshots/dashboard_main.png"
    if os.path.exists(img7_path):
        s7.shapes.add_picture(img7_path, Inches(7.8), Inches(2.0), width=Inches(5.0))

    # ── SLIDE 8: Results and Discussion ─────────────────────────────────────
    s8 = prs.slides[7]
    tf8 = get_or_create_body_textbox(s8)
    tf8.clear()
    bullets8 = [
        ("Sub-3 Second Emergency Latency", "Fall detection (MPU6050 freefall signature) & SOS triggers fire dashboard alerts and Telegram notifications in <2.8 seconds."),
        ("5km LoRa Range", "Successfully tested 433MHz LoRa packet delivery across rural non-line-of-sight environments without packet drop."),
        ("High Medical Accuracy", "MAX30102 pulse oximeter & MLX90614 temperature sensors show <1.5% variance compared to clinical grade monitors."),
        ("100% System Uptime", "Auto-switching mock simulation fallback ensures zero dashboard downtime even during network disconnects.")
    ]
    for idx, (title, desc) in enumerate(bullets8):
        p = tf8.paragraphs[0] if idx == 0 else tf8.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_SLATE

    img8_path = "/home/siddhartha/Desktop/copied/backfol/AyuLink/ElderCare/photos/fids_detection_alert.jpg"
    if os.path.exists(img8_path):
        s8.shapes.add_picture(img8_path, Inches(7.8), Inches(2.0), width=Inches(5.0))

    # ── SLIDE 9: Impact on Society ───────────────────────────────────────────
    s9 = prs.slides[8]
    tf9 = get_or_create_body_textbox(s9, width=Inches(11.5))
    tf9.clear()
    bullets9 = [
        ("Saving Lives in the Golden Window", "Reduces rural emergency response time from 47 minutes toward 8 minutes for cardiac and fall events."),
        ("Universal Healthcare Access", "Extends remote monitoring to 600M+ underserved rural Indians and 30,000+ Primary Health Centers."),
        ("Economic Accessibility", "Lowers annual patient monitoring cost by 96% (from ₹6,000+/yr to ₹216/yr)."),
        ("Empowerment for Non-English Speakers", "Telugu/Hindi voice notifications allow uneducated family members and ASHA field workers to take instant action.")
    ]
    for idx, (title, desc) in enumerate(bullets9):
        p = tf9.paragraphs[0] if idx == 0 else tf9.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE

    # ── SLIDE 10: Future Scope ───────────────────────────────────────────────
    s10 = prs.slides[9]
    tf10 = get_or_create_body_textbox(s10, width=Inches(11.5))
    tf10.clear()
    bullets10 = [
        ("District-Wide LoRa Mesh", "Chain multi-village LoRa repeater gateways for full district-level network coverage."),
        ("Edge AI Triage", "Deploy lightweight TFLite Micro models directly on ESP32-S3 microcontrollers for zero-latency local triage."),
        ("ABHA & Government PHC Integration", "Direct API integration with India’s Ayushman Bharat Digital Health Mission."),
        ("Custom Waterproof Hardware v2", "Mass-manufactured IP67 waterproof wristband PCB with 7-day battery life and solar charging nodes.")
    ]
    for idx, (title, desc) in enumerate(bullets10):
        p = tf10.paragraphs[0] if idx == 0 else tf10.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE

    # ── SLIDE 11: Conclusions ────────────────────────────────────────────────
    s11 = prs.slides[10]
    tf11 = get_or_create_body_textbox(s11, width=Inches(11.5))
    tf11.clear()
    bullets11 = [
        ("Democratizing Healthcare", "AyuLink proves that life-saving emergency medical technology does not require expensive smartphones or 4G infrastructure."),
        ("End-to-End Execution", "Successfully combined physical IoT hardware, LoRa radio mesh, FastAPI backend, LLaMA 3.1 AI intelligence, and Next.js 16 Web Dashboard into a functional product."),
        ("Scalable Impact", "Solves the rural healthcare gap for ₹18/month while remaining fully scalable for urban homecare and digital telemedicine.")
    ]
    for idx, (title, desc) in enumerate(bullets11):
        p = tf11.paragraphs[0] if idx == 0 else tf11.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE

    # ── SLIDE 12: References ─────────────────────────────────────────────────
    s12 = prs.slides[11]
    tf12 = get_or_create_body_textbox(s12, width=Inches(11.5))
    tf12.clear()
    refs = [
        "1. World Health Organization (WHO), 'Global Report on Health Emergencies and Rural Access,' 2024.",
        "2. Ministry of Health and Family Welfare, Govt. of India, 'Ayushman Bharat Digital Health Mission Guidelines,' 2025.",
        "3. Semtech Corporation, 'SX1276/77/78/79 Long Range, Low Power Transceiver Data Sheet,' 2023.",
        "4. Groq Inc., 'Llama 3.1 Ultra-Low Latency Inference Engine Documentation,' 2025.",
        "5. Next.js & React Core Team, 'Next.js 16 App Router & Server Components Architecture,' 2026."
    ]
    for idx, ref in enumerate(refs):
        p = tf12.paragraphs[0] if idx == 0 else tf12.add_paragraph()
        p.text = ref
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE

    # ── SLIDE 13: Thank You ──────────────────────────────────────────────────
    s13 = prs.slides[12]
    for shape in s13.shapes:
        if shape.has_text_frame and 'TextBox 8' in shape.name:
            tf13 = shape.text_frame
            tf13.clear()
            p0 = tf13.paragraphs[0]
            p0.text = "Thank You!"
            p0.font.bold = True
            p0.font.size = Pt(36)
            p0.font.color.rgb = TEAL

            p1 = tf13.add_paragraph()
            p1.text = "AyuLink — Smart Health. Zero Boundaries."
            p1.font.bold = True
            p1.font.size = Pt(20)
            p1.font.color.rgb = DARK_SLATE

            p2 = tf13.add_paragraph()
            p2.text = "\nTeam FightClub | TechSphere 2026 (Woxsen University)\nPresenters: Siddhartha & Anirudh\nLive Dashboard: ayulink-woad.vercel.app\nGitHub: github.com/siddharthathula/AyuLink"
            p2.font.size = Pt(14)
            p2.font.color.rgb = MUTED_SLATE

    # Save outputs
    prs.save(OUTPUT_PATH_1)
    prs.save(OUTPUT_PATH_2)
    print(f"Successfully populated and saved presentation to:\n - {OUTPUT_PATH_1}\n - {OUTPUT_PATH_2}")

if __name__ == '__main__':
    populate_presentation()
